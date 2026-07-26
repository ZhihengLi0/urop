#!/usr/bin/env python3
"""Build one standalone backup gallery PDF per zip: backup_zip{N}.pdf.

Same treatment as the Z7 gallery inside the main deck: every result figure
of the zip, in processing order, split at panel boundaries into readable
pages (channel figures: two full-width strips per page; event x channel
grids: 3 event rows x half the channels per page). Figure headers with the
stamped processing text are dropped — only plots, titles, axes and legends
remain.

Sources (MSI):
  lp_fit_align/results/plots/<kind>/zip{N}_*.png
  deliverables/nxm/plots/zip{N}_pca_templates.png
  deliverables/1x1/plots/{PT,PS1,PS2}/zip{N}_*.png
  presentation/figures/kline_zip{N}.png

Usage: python3 build_backup_pdfs.py --zips 1 4 6 ...   (pptx written to a
temp dir; convert to PDF with LibreOffice afterwards)
"""
import argparse
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.abspath(os.path.join(HERE, ".."))
RES = os.path.join(BASE, "lp_fit_align", "results", "plots")
DELIV = os.path.join(BASE, "deliverables")
OUT_DIR = os.path.join(HERE, "backup_galleries")
TMP_DIR = os.path.join(OUT_DIR, "_pptx")
SPLIT_DIR = os.path.join(OUT_DIR, "_split")
for d in (OUT_DIR, TMP_DIR, SPLIT_DIR):
    os.makedirs(d, exist_ok=True)

NAVY = RGBColor(0x1F, 0x38, 0x64)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x77, 0x77, 0x77)
IN = Inches
SW, SH = Inches(13.333), Inches(7.5)


# ---------------------------------------------------------- panel slicing
def _white_rows(img):
    import numpy as np
    g = np.asarray(img.convert("L"))
    return (g.min(axis=1) > 247)


def _white_cols(img):
    import numpy as np
    g = np.asarray(img.convert("L"))
    return (g.min(axis=0) > 247)


def snap_cut(white, target, window):
    lo = max(1, target - window)
    hi = min(len(white) - 1, target + window)
    best, bestd = None, None
    for y in range(lo, hi):
        if white[y]:
            d = abs(y - target)
            if bestd is None or d < bestd:
                best, bestd = y, d
    return best if best is not None else target


def _content_blocks(mask):
    blocks, start = [], None
    for i, w in enumerate(mask):
        if not w and start is None:
            start = i
        elif w and start is not None:
            blocks.append((start, i)); start = None
    if start is not None:
        blocks.append((start, len(mask)))
    return blocks


def _panel_units(img):
    """(start, end) per panel: panel box + its hugging title above and the
    tick/label text below; the header text above the first panel is dropped."""
    white = _white_rows(img)
    blocks = _content_blocks(white)
    hmax = max(b[1] - b[0] for b in blocks)
    thr = max(60, int(hmax * 0.35))
    panels = [b for b in blocks if b[1] - b[0] >= thr]
    small = [b for b in blocks if b[1] - b[0] < thr]
    starts = []
    for ps, pe in panels:
        above = sorted(b for b in small if b[1] <= ps + 2)
        top = ps
        if above and ps - above[-1][1] <= 12:      # the title hugs its panel
            top = above[-1][0]
            for b in reversed(above[:-1]):         # merge artifact slivers
                if top - b[1] <= 6:
                    top = b[0]
                else:
                    break
        starts.append(max(0, top - 4))
    units, tails = [], []
    for i, (ps, pe) in enumerate(panels[:-1]):
        end = starts[i + 1] - 4
        units.append((starts[i], end))
        tails.append(end - pe)
    ps, pe = panels[-1]
    tail = sorted(tails)[len(tails) // 2] if tails else 60
    units.append((starts[-1], min(img.size[1], pe + tail)))
    return units


# ---------------------------------------------------------- deck helpers
def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def title(s, text, sub=None):
    tb = textbox(s, IN(0.55), IN(0.28), IN(12.3), IN(0.75))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.color.rgb = Pt(27), True, NAVY
    r.font.name = "Arial"
    if sub:
        p2 = tb.text_frame.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size, r2.font.color.rgb, r2.font.name = Pt(14), GRAY, "Arial"
    ln = s.shapes.add_shape(1, IN(0.55), IN(1.02), IN(12.3), Emu(18000))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def _gal_slide(prs, head, sub_txt, info):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, head, sub=sub_txt)
    tb = textbox(s, IN(0.55), IN(1.18), IN(12.3), IN(0.62))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = info
    r.font.size, r.font.color.rgb, r.font.name = Pt(11.5), DARK, "Arial"
    return s


def _place_stack(s, crops, l, t, box_w, box_h, gap=Inches(0.22)):
    n = len(crops)
    cell_h = (box_h - gap * (n - 1)) / n
    for i, (cp, (cw, ch)) in enumerate(crops):
        sc = min(box_w / cw, cell_h / ch)
        dw, dh = int(cw * sc), int(ch * sc)
        left = l + int((box_w - dw) / 2)
        top = t + i * (cell_h + gap) + int((cell_h - dh) / 2)
        s.shapes.add_picture(cp, left, top, width=dw, height=dh)


def pic(s, path, l, t, max_w, max_h, caption=None):
    w0, h0 = Image.open(path).size
    scale = min(max_w / w0, max_h / h0)
    w, h = int(w0 * scale), int(h0 * scale)
    left = l + int((max_w - w) / 2)
    s.shapes.add_picture(path, left, t, width=w, height=h)
    if caption:
        tb = textbox(s, l, t + h + Emu(30000), max_w, Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = caption
        r.font.size, r.font.italic, r.font.color.rgb = Pt(11), True, GRAY
        r.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER


def channel_fig_pages(prs, z, path, short, info, per_page=2):
    img = Image.open(path)
    w0 = img.size[0]
    units = _panel_units(img)
    n = len(units)
    base = f"z{z}_" + os.path.splitext(os.path.basename(path))[0]
    fname = os.path.basename(path)
    npages = (n + per_page - 1) // per_page
    for p in range(npages):
        i0, i1 = p * per_page, min((p + 1) * per_page, n)
        crops = []
        for i in range(i0, i1):
            crop = img.crop((0, units[i][0], w0, units[i][1]))
            cp = os.path.join(SPLIT_DIR, f"{base}_r{i+1}.png")
            crop.save(cp)
            crops.append((cp, crop.size))
        s = _gal_slide(prs, f"Z{z} · {short}  ({p+1}/{npages})",
                       f"figure: {fname} — channel panels {i0+1}–{i1} of {n}, "
                       f"top to bottom", info)
        _place_stack(s, crops, IN(0.55), IN(1.95), IN(12.3), IN(5.3))
        notes(s, f"Z{z} gallery, {short}, page {p+1} of {npages}. {info}")


def grid_fig_pages(prs, z, path, short, info, rows_per_page=3):
    img = Image.open(path)
    w0 = img.size[0]
    units = _panel_units(img)
    n = len(units)
    body = img.crop((0, units[0][0], w0, units[-1][1]))
    xh = snap_cut(_white_cols(body), w0 // 2, w0 // 10)
    base = f"z{z}_" + os.path.splitext(os.path.basename(path))[0]
    fname = os.path.basename(path)
    groups = [(a, min(a + rows_per_page, n))
              for a in range(0, n, rows_per_page)]
    pages = [(a, b, side) for a, b in groups for side in (0, 1)]
    for p, (a, b, side) in enumerate(pages):
        x0, x1 = (0, xh) if side == 0 else (xh, w0)
        half = "left" if side == 0 else "right"
        crop = img.crop((x0, units[a][0], x1, units[b - 1][1]))
        cp = os.path.join(SPLIT_DIR, f"{base}_e{a+1}-{b}_{half}.png")
        crop.save(cp)
        s = _gal_slide(prs, f"Z{z} · {short}  ({p+1}/{len(pages)})",
                       f"figure: {fname} — event rows {a+1}–{b} of {n}, "
                       f"{half} half of the channels", info)
        _place_stack(s, [(cp, crop.size)], IN(0.55), IN(1.95), IN(12.3), IN(5.3))
        notes(s, f"Z{z} gallery, {short}, page {p+1} of {len(pages)}. {info}")


# ---------------------------------------------------------- per-zip build
def build_zip(z):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # intro
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, f"Backup — Z{z} full results gallery",
          sub="every diagnostic and result figure of the pipeline, in processing order")
    tb = textbox(s, IN(0.55), IN(1.6), IN(12.3), IN(4.6))
    for i, line in enumerate([
            "Dataset: the K-line selection window",
            "Example grids: LP + fit and raw vs fit, event × channel",
            "Fit diagnostics: NRMSE, fitted pretrigger, τ_rise / τ_fall distributions",
            "Aligned fans: no cut → NRMSE ≤ 0.4 → + τ_rise ≤ 0.3 ms; kept vs rejected; measured traces",
            "Special populations: NRMSE-rejected events, well-fit slow-rise events",
            "Templates: PCA nxm0–4 per channel, summed PT / PS1 / PS2",
    ]):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        r = p.add_run(); r.text = "• " + line
        r.font.size, r.font.color.rgb, r.font.name = Pt(16), DARK, "Arial"
        p.space_after = Pt(10)
    notes(s, f"Standalone Z{z} gallery: every figure the pipeline produced.")

    # dataset window
    s = _gal_slide(prs, f"Z{z} · K-line selection window",
                   f"figure: kline_zip{z}.png",
                   f"Summed PTOFamps spectrum of Z{z}; red line = Prof. Saab's 10.37 keV "
                   "K-line position; the ×/÷ 1.35 window around it is the event selection; "
                   "everything inside is cached raw.")
    pic(s, os.path.join(HERE, "figures", f"kline_zip{z}.png"),
        IN(2.65), IN(2.0), IN(8.0), IN(5.2),
        caption=f"Z{z} PTOFamps spectrum with the K-line selection window")

    grid_fig_pages(prs, z, f"{RES}/fit_examples/zip{z}_fit_examples.png",
        "fit_examples — LP trace vs 2-exp fit",
        "First 15 cached events × all channels; blue = 100 kHz LP trace, red = 2-exp fit. "
        "A real event fits well in every channel at once; a noise trigger fails across the whole row.")
    grid_fig_pages(prs, z, f"{RES}/raw_vs_fit_examples/zip{z}_raw_vs_fit_examples.png",
        "raw_vs_fit_examples — raw trace vs fit",
        "Same 15 events × all channels; gray = raw unfiltered trace, blue = LP, red = fit, "
        "per-panel NRMSE top-right. Noise vs pulse size is directly visible in the raw data.")

    channel_fig_pages(prs, z, f"{RES}/nrmse/zip{z}_nrmse.png",
        "nrmse — NRMSE distribution per channel",
        "NRMSE = RMS(fit residual)/fitted peak, log-log. Bimodal: good fits vs noise "
        "triggers; the valley at ≈0.4 sets quality cut 1.")
    channel_fig_pages(prs, z, f"{RES}/pretrigger/zip{z}_pretrigger.png",
        "pretrigger — fitted onset per channel",
        "The onset is a FREE fit parameter; pinning it inside the fit would bias every "
        "other parameter.")
    channel_fig_pages(prs, z, f"{RES}/time_constants/zip{z}_time_constants.png",
        "time_constants — τ_rise / τ_fall per channel",
        "Fast-pulse population at τ_rise ≈ 0.1 ms; isolated spikes at the parameter bounds "
        "are noise fits pinned at the fit limits, removed by the NRMSE cut.")

    channel_fig_pages(prs, z, f"{RES}/fitted_curves_overlay/zip{z}_fitted_curves_overlay.png",
        "fitted_curves_overlay — fan, no cut",
        "Every fit_ok fitted curve at the common pretrigger 16050, peak-normalized. Pure "
        "shape distribution; the broad slow curves are noise fits that the cuts remove.")
    channel_fig_pages(prs, z, f"{RES}/fitted_curves_overlay/zip{z}_fitted_curves_overlay_nrmse0.4.png",
        "fitted_curves_overlay — after NRMSE ≤ 0.4",
        "Same fan after quality cut 1: the noise fan collapses, the fast-pulse bundle remains.")
    channel_fig_pages(prs, z, f"{RES}/fitted_curves_overlay/zip{z}_fitted_curves_overlay_nrmse0.4_trise0.30ms.png",
        "fitted_curves_overlay — after both cuts",
        "After NRMSE ≤ 0.4 AND τ_rise ≤ 0.3 ms — the exact PCA input population: the clean "
        "fast-pulse bundle only.")
    channel_fig_pages(prs, z, f"{RES}/overlay_fan_cut/zip{z}_overlay_fan_cut_nrmse0.4.png",
        "overlay_fan_cut — kept vs rejected fits",
        "Gray-blue = aligned measured traces; green = fits with NRMSE ≤ 0.4 (kept); red = "
        "NRMSE > 0.4 (cut away). Median NRMSE and counts of both populations stamped per channel.")
    channel_fig_pages(prs, z, f"{RES}/aligned_overlay/zip{z}_lp_aligned_overlay.png",
        "aligned_overlay — aligned measured traces",
        "Blue = measured LP traces shifted by (fitted pretrigger − 16050); red = mean of all "
        "fit_ok events; orange = NRMSE-weighted mean of the fitted curves. Rise edges line up at 16050.")

    grid_fig_pages(prs, z, f"{RES}/slow_rise_events/zip{z}_slow_rise_events.png",
        "slow_rise_events — NRMSE-rejected events",
        "Events with median NRMSE > 0.4 across channels: the raw traces show no pulse in "
        "any channel — the rejected population is noise triggers, not physics.")
    shadow = f"{RES}/shadow_events/zip{z}_shadow_events.png"
    if os.path.exists(shadow):
        grid_fig_pages(prs, z, shadow,
            "shadow_events — well-fit slow-rise events",
            "Median NRMSE ≤ 0.4 AND median τ_rise > 0.2 ms: genuinely slow, well-fit pulses "
            "— kept, and exactly what the NxM multi-templates are for.")
    slow_fall = f"{RES}/slow_fall_events/zip{z}_PDS2_slow_fall.png"
    if os.path.exists(slow_fall):
        grid_fig_pages(prs, z, slow_fall,
            "slow_fall_events — long-τ_fall study (PDS2)",
            "Random events with τ_fall > 1.5 ms drawn in every channel: the long fall is the "
            "PDS2 low-frequency hardware noise, not slow physics.")
    trise_rm = f"{RES}/trise_removed/zip{z}_trise_removed.png"
    if os.path.exists(trise_rm):
        grid_fig_pages(prs, z, trise_rm,
            "trise_removed — noise that NRMSE missed",
            "Events that PASS NRMSE ≤ 0.4 but are removed by τ_rise ≤ 0.3 ms: a slow 2-exp "
            "hugging a slow baseline drift — no clean fast pulse in the raw trace.")

    channel_fig_pages(prs, z, f"{DELIV}/nxm/plots/zip{z}_pca_templates.png",
        "deliverables/nxm — PCA templates nxm0–4",
        "nxm0 = mean curve; nxm1–4 = PCA components (oscillating basis vectors, may be "
        "negative), all peak-normalized; population = fit_ok + NRMSE ≤ 0.4 + τ_rise ≤ 0.3 ms.")

    # summed 1x1 templates
    s = _gal_slide(prs, f"Z{z} · 1x1 summed templates PT / PS1 / PS2",
                   f"figures: deliverables/1x1/plots/{{PT,PS1,PS2}}/zip{z}_*.png",
                   "Single-template (1x1) summed curves: peak-normalized average of the "
                   "per-channel nxm0 templates — PT = all channels, PS1 / PS2 = side-1 / side-2 only.")
    for k, name in enumerate(["PT", "PS1", "PS2"]):
        path = f"{DELIV}/1x1/plots/{name}/zip{z}_{name}.png"
        img = Image.open(path)
        u = _panel_units(img)
        crop = img.crop((0, u[0][0], img.size[0], u[-1][1]))
        cp = os.path.join(SPLIT_DIR, f"z{z}_{name}_clean.png")
        crop.save(cp)
        pic(s, cp, IN(0.4 + 4.3 * k), IN(2.3), IN(4.1), IN(4.2),
            caption=f"Z{z} {name}")

    # page numbers
    slides = list(prs.slides)
    total = len(slides)
    for i, s in enumerate(slides, start=1):
        pn = s.shapes.add_textbox(SW - Inches(1.35), SH - Inches(0.45),
                                  Inches(1.1), Inches(0.3))
        p = pn.text_frame.paragraphs[0]
        r = p.add_run(); r.text = f"{i} / {total}"
        r.font.size, r.font.color.rgb, r.font.name = Pt(11), GRAY, "Arial"
        p.alignment = PP_ALIGN.RIGHT

    out = os.path.join(TMP_DIR, f"backup_zip{z}.pptx")
    prs.save(out)
    print(f"zip{z}: {total} slides -> {out}")
    return out


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--zips", type=int, nargs="+", required=True)
    args = ap.parse_args()
    for z in args.zips:
        build_zip(z)
