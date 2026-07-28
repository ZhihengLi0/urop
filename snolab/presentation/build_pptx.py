#!/usr/bin/env python3
"""Build the 10-min SNOLAB R4 template-project slide deck (16:9, English)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "SNOLAB_R4_pulse_templates_20260715.pptx")

NAVY = RGBColor(0x1F, 0x38, 0x64)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0xC0, 0x50, 0x4D)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def title(s, text, sub=None):
    tb = textbox(s, Inches(0.55), Inches(0.28), Inches(12.3), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.color.rgb = Pt(27), True, NAVY
    r.font.name = "Arial"
    if sub:
        p2 = tb.text_frame.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size, r2.font.color.rgb, r2.font.name = Pt(14), GRAY, "Arial"
    # thin rule under the title
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.02), Inches(12.3), Emu(18000))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()


def bullets(s, items, l, t, w, h, size=15):
    """items: list of (level, text) or str (level 0)."""
    tb = textbox(s, l, t, w, h)
    tf = tb.text_frame
    first = True
    for it in items:
        lvl, txt = it if isinstance(it, tuple) else (0, it)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = DARK if lvl == 0 else GRAY
        r.font.name = "Arial"
    return tb


def pic(s, name, l, t, max_w, max_h, caption=None, stretch=False):
    path = os.path.join(FIG, name)
    w0, h0 = Image.open(path).size
    if stretch:                      # fill the whole box, aspect not preserved
        w, h = int(max_w), int(max_h)
    else:
        scale = min(max_w / w0, max_h / h0)
        w, h = int(w0 * scale), int(h0 * scale)
    # center inside the box
    left = l + int((max_w - w) / 2)
    s.shapes.add_picture(path, left, t, width=w, height=h)
    if caption:
        tb = textbox(s, l, t + h + Emu(30000), max_w, Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = caption
        r.font.size, r.font.italic, r.font.color.rgb = Pt(11), True, GRAY
        r.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
    return h


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


IN = Inches  # shorthand

# ---------------------------------------------------------------- 1 · title
s = slide()
tb = textbox(s, IN(0.9), IN(2.3), IN(11.5), IN(1.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "NxM Templates for SuperCDMS SNOLAB Run 4 using Ge Activation Data"
r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(36), True, NAVY, "Arial"
tb2 = textbox(s, IN(0.9), IN(3.7), IN(11.5), IN(0.9))
p = tb2.text_frame.paragraphs[0]
r = p.add_run()
r.text = "K-line event selection  ·  1x1 and NxM (PCA) templates"
r.font.size, r.font.color.rgb, r.font.name = Pt(18), GRAY, "Arial"
tb3 = textbox(s, IN(0.9), IN(5.6), IN(11.5), IN(0.8))
p = tb3.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Zhiheng Li  —  July 2026"
r.font.size, r.font.color.rgb, r.font.name = Pt(16), DARK, "Arial"
notes(s, "Today I will present the phonon pulse-template work for SNOLAB Run 4: "
         "how we selected events on the Ge-activation K-line, the fit-and-align "
         "pipeline, the quality cuts and how each was derived, and the two "
         "template families we delivered for all 13 detectors.")


# ------------------------------------------------- 1b · table of contents
s = slide()
title(s, "Table of Contents")
_toc = textbox(s, IN(1.1), IN(1.7), IN(11.0), IN(4.9))
for _i, (_head, _rest) in enumerate([
        ("Dataset", "Ge activation K-line event selection, raw trace cache (13 detectors, 23–30 series each, ≈1–2 h per series)"),
        ("Pulse fit", "low-pass → fit → align"),
        ("Data cleaning", "fit_ok, NRMSE valley cut (0.4), rise-time ceiling (0.3 ms)"),
        ("Cross-checks with real pulses", "cuts remove noise, real pulses stay; slow-fall follow-up"),
        ("Template generation", "1x1 weighted template, NxM PCA templates"),
        ("Template file and future steps", "cdmsbats ROOT files, ready for merge; three-, four-exponential fits"),
        ("Backup", "weak detectors, rise-time cut details, full results for every zip")]):
    _p = _toc.text_frame.paragraphs[0] if _i == 0 else _toc.text_frame.add_paragraph()
    _p.space_after = Pt(10)
    _r = _p.add_run(); _r.text = f"{_i+1}.  {_head}: "
    _r.font.size, _r.font.bold, _r.font.color.rgb, _r.font.name = Pt(19), True, NAVY, "Arial"
    _r = _p.add_run(); _r.text = _rest
    _r.font.size, _r.font.color.rgb, _r.font.name = Pt(18), DARK, "Arial"
notes(s, "Quick roadmap: the dataset and how events were selected, the fit and "
         "alignment, the two cleaning cuts and the checks behind them, the two "
         "template families we deliver, and the future steps. Backup slides "
         "hold the weak-detector picture and the full results for every zip.")

# ------------------------------------------------- 2 · motivation / selection
s = slide()
title(s, "Ge activation dataset")
_tb = textbox(s, IN(0.25), IN(1.14), IN(12.6), IN(1.85))
_tf = _tb.text_frame
def _row(first, parts, size=13.5):
    _p = _tf.paragraphs[0] if first else _tf.add_paragraph()
    _p.space_after = Pt(4)
    for _txt, _c, _b in parts:
        _r = _p.add_run(); _r.text = _txt
        _r.font.size, _r.font.name = Pt(size), "Arial"
        _r.font.color.rgb = _c; _r.font.bold = _b
_row(True, [("• Per detector, Prof. Saab's study marked the 10.37 keV K-line position in ", DARK, False),
            ("PTOFamps", ACCENT, True),
            (" (the red line in each panel below)", DARK, False)])
_row(False, [("• A ", DARK, False), ("PTOFamps", ACCENT, True),
             (" window bracketing the red line: position ×/÷ 1.35, a rough eyeballed choice; deliberately minimal, no other cuts", DARK, False)])
_row(False, [("• For every selected event: cache the fully unprocessed raw MIDAS traces of all channels + event metadata, "
              "saved as per-series pickle files (reusable for any later analysis)", DARK, False)])
_row(False, [("• Events in the window per zip:  Z1 10.1k   Z4 26.1k   Z6 17.0k   Z7 2.2k   Z9 2.6k   Z10 7.6k   Z13 9.1k   "
              "Z15 4.2k   Z16 1.6k   Z18 22.1k   Z19 23.3k   Z22 25.0k   Z24 22.9k   (total 174k)", DARK, False)], size=12.5)
_row(False, [("• Focus: Z7, the best detector; examples use its PBS1 channel; other detectors in the backup slides", NAVY, True)])
# credit box, top-right of the title band
ctb = textbox(s, IN(9.35), IN(0.30), IN(3.5), IN(0.72))
_p = ctb.text_frame.paragraphs[0]
_r = _p.add_run(); _r.text = "Credit: Prof. Tarek Saab"
_r.font.size, _r.font.color.rgb, _r.font.name = Pt(10), GRAY, "Arial"
_r.hyperlink.address = "https://confluence.slac.stanford.edu/spaces/CDMS/pages/716899864/Ge+Activation+Data+-+Ops+Shift+260612"
_p.alignment = PP_ALIGN.RIGHT
pic(s, "kline_all_zips.png", IN(0.42), IN(3.05), IN(12.5), IN(4.05),
    "All 13 detectors, PTOFamps spectra (Prof. Saab's study)")
notes(s, "This is the Ge activation dataset from Run 4, credit to Professor "
         "Saab's study on the SuperCDMS Confluence. The goal of the whole "
         "project is one clean, minimally biased pulse population per "
         "detector and channel. After Cf activation every detector shows the "
         "10.37 keV K-line; the study located its position in the PTOFamps "
         "spectrum of each detector, the red line in every panel. Our event "
         "selection is exactly one condition: a window around the red line, "
         "a factor 1.35 both ways, and nothing else. Every selected event is "
         "cached with its fully unprocessed raw MIDAS traces and metadata as "
         "per-series pickle files, so the sample can be reused for any later "
         "analysis. The window holds from about 1.6 thousand events on Z16 "
         "up to 26 thousand on Z4, 174 thousand in total; Z7 has about 2.2 "
         "thousand. In this talk I focus on Z7, the best detector; the other "
         "detectors are covered in the backup slides.")

# ------------------------------------------------------ 3 · fit quality grids (overview)
s = slide()
title(s, "Fit examples")
bullets(s, [
    "Real events and noise triggers can be told apart by visual comparison (no cut applied at this stage)",
], IN(0.55), IN(1.2), IN(12.3), IN(0.7), size=14)
def _biglabel(x, w, txt, color):
    _tb = textbox(s, x, IN(1.95), w, IN(0.5))
    _p = _tb.text_frame.paragraphs[0]
    _r = _p.add_run(); _r.text = txt
    _r.font.size, _r.font.bold, _r.font.color.rgb, _r.font.name = Pt(24), True, color, "Arial"
    _p.alignment = PP_ALIGN.CENTER
_biglabel(IN(0.2), IN(6.55), "NOISE", ACCENT)
_biglabel(IN(6.6), IN(6.55), "REAL PULSE", NAVY)
pic(s, "fit_examples_zip7_noise.png", IN(0.2), IN(2.55), IN(6.55), IN(4.0))
pic(s, "fit_examples_zip7_good.png", IN(6.6), IN(2.55), IN(6.55), IN(4.0))
notes(s, "Before the algorithm, a quick look at the data itself. We use "
         "event-by-channel grids: each row is one event, each column one "
         "channel, with the low-passed trace in blue and a fitted pulse model "
         "in red on top. By visual comparison you can already tell the two "
         "kinds apart: on the left, noise triggers; on the right, K-line "
         "events. No cut is applied at this stage, and an event is not "
         "necessarily good or bad in every channel at once. "
         "How we actually fit and align each trace is the next slide, and how "
         "we turn this into a quantitative cut comes shortly after.")

# ------------------------------------------------------ 4 · per-trace algorithm
s = slide()
title(s, "Per-trace algorithm: low-pass → fit → align")
def steps_box(top, height, items):
    _tb = textbox(s, IN(0.55), top, IN(12.3), height)
    _first = True
    for head, rest in items:
        p = _tb.text_frame.paragraphs[0] if _first else _tb.text_frame.add_paragraph()
        _first = False
        p.space_after = Pt(6)
        r = p.add_run(); r.text = head
        r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(16), True, NAVY, "Arial"
        r = p.add_run(); r.text = rest
        r.font.size, r.font.color.rgb, r.font.name = Pt(16), DARK, "Arial"

steps_box(IN(1.2), IN(1.35), [
    ("1.  Low-pass", ": 100 kHz Butterworth"),
    ("2.  Normalize", ": subtract the baseline (median of the pre-pulse samples), scale the trace to unit peak"),
    ("3.  Fit", ": two-exponential model:"),
])
pic(s, "formula_2exp.png", IN(2.6), IN(2.62), IN(8.1), IN(0.95))
# fit parameter definitions, right under the formula
_pdef = textbox(s, IN(1.15), IN(3.58), IN(11.6), IN(0.35))
_pp = _pdef.text_frame.paragraphs[0]
_pr = _pp.add_run()
_pr.text = ("Fit parameters:  A amplitude   ·   τ_rise rise time constant   ·   "
            "τ_fall fall time constant   ·   t₀ pulse onset (pretrigger)   ·   b baseline")
_pr.font.size, _pr.font.italic = Pt(12.5), True
_pr.font.color.rgb, _pr.font.name = GRAY, "Arial"
steps_box(IN(3.95), IN(0.55), [
    ("4.  Two quality checks per trace", " (calculated and recorded here, not cut yet):"),
])
# the two quality-number sub-lines, plain English
q_tb = textbox(s, IN(1.15), IN(4.42), IN(11.6), IN(0.9))
for i, (lead, rest) in enumerate([
        ("fit_ok", ": amplitude is positive, and the pulse rises faster than it falls"),
        ("NRMSE", ": RMS of the fit residual ÷ pulse peak  (smaller = better fit)")]):
    p = q_tb.text_frame.paragraphs[0] if i == 0 else q_tb.text_frame.add_paragraph()
    p.space_after = Pt(4)
    r = p.add_run(); r.text = "•  " + lead
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(15), True, DARK, "Arial"
    r = p.add_run(); r.text = rest
    r.font.size, r.font.color.rgb, r.font.name = Pt(15), DARK, "Arial"
notes(s, "The per-trace algorithm. First a low-pass filter to "
         "smooth away the high-frequency noise, then the baseline is "
         "subtracted and the trace is scaled to unit peak - without that "
         "normalization the traces could not be overlaid or compared. The "
         "core is the fit: a two-exponential function, one exponential for "
         "the rise and one for the fall, with all five parameters free, "
         "including the pretrigger - the pulse start time. The pretrigger is "
         "never pinned: real trigger times vary, and the fitted pretriggers "
         "cluster about 230 samples after the nominal value - pinning it "
         "would bias the time constants. Each fit records two quality "
         "numbers: fit_ok, a physicality check - positive amplitude, rise "
         "faster than fall - and the NRMSE, the normalized root-mean-square "
         "error, the fit residual divided by the pulse height. At this stage "
         "we only record them. Step five, alignment, and its output get "
         "their own slide next. Q and A backup: why fit_ok demands rise "
         "faster than fall - the rise time is set by L over R, an electrical "
         "property of the readout circuit, while the fall time is set by C "
         "over G, a thermal property of the detector; physically the "
         "electrical rise is the faster of the two.")

# ------------------------------------------------------ 5 · alignment result
s = slide()
title(s, "5. Align")
pic(s, "fan_zip7_PBS1_allfits.png", IN(1.7), IN(1.45), IN(9.9), IN(5.5),
    "Z7 PBS1: every fitted event drawn at the common pretrigger, no cuts of any kind (n = 2203, inverted fits included)")
notes(s, "Step five: alignment. Each measured trace is shifted by its fitted "
         "pretrigger minus the nominal 16050 - a pure translation, nothing "
         "about the waveform is regenerated. Drawn at the common pretrigger, "
         "the fitted curves stack into one shape family. This is every "
         "fitted event on this channel - two thousand and eight curves, no "
         "fit_ok requirement, no NRMSE, no cuts of any kind, because at this "
         "point in the story we have not selected anything yet. You can see "
         "the dominant fast bundle and the stragglers spreading off it; how "
         "we separate them quantitatively is the next slide.")

# ------------------------------------------------------ 6 · NRMSE cut
# ------------------------------------------------- 6b · quality cut 1: fit_ok
s = slide()
title(s, "Quality cut 1: fit_ok")
bullets(s, [
    "fit_ok: amplitude > 0 AND rise faster than fall (the fit is a physical pulse)",
    "Z7 PBS1: 2008 of 2206 events pass; 195 removed (red, inverted or negative shapes), 3 fits did not converge",
], IN(0.55), IN(1.2), IN(12.3), IN(1.1), size=15)
pic(s, "fan_zip7_PBS1_allfits.png", IN(0.3), IN(2.45), IN(6.35), IN(4.4),
    "All fitted events, no cuts (n = 2203 drawn, inverted fits included)")
pic(s, "fan_zip7_PBS1_fitok_cut.png", IN(6.65), IN(2.45), IN(6.35), IN(4.4),
    "fit_ok: kept (blue) vs removed (red, drawn at |peak| normalization)")
notes(s, "The first quality cut is fit_ok, a pure physics check: the "
         "amplitude must be positive and the rise must be faster than the "
         "fall. On the left the same all-events fan; on the right the fits "
         "removed by fit_ok are drawn in red, inverted or negative shapes, "
         "clearly non-physical. On PBS1 it removes 195 of 2206 events and "
         "three more fits never converged; 2008 curves move on.")

s = slide()
title(s, "Quality cut 2: NRMSE ≤ 0.4 — where the number comes from")
bullets(s, [
    "NRMSE of fit_ok events is bimodal: good fits (median ≈ 0.05–0.1) vs noise triggers (≈ 1–2), valley at ≈ 0.4–0.5",
    "The cut sits in the valley: eyeball read off the distribution",
    "Worse detectors (Z1, Z4, Z6, Z18, Z19, Z22, Z24): same bimodal picture, noise peak dominates",
], IN(0.55), IN(1.2), IN(12.3), IN(1.55), size=14)
pic(s, "nrmse_zip7_PBS1.png", IN(1.4), IN(3.1), IN(10.5), IN(3.4),
    "Z7 PBS1: two clean populations, log-log axes")
notes(s, "First quality cut. The NRMSE distribution of physical fits is "
         "bimodal on every detector: a good-fit population around 0.05 to "
         "0.1, and a noise-trigger population around 1 to 2, separated by a "
         "valley at about 0.4. We place the cut in the valley - it is read "
         "off the distribution itself. This is Z7 PBS1, the channel used "
         "throughout the talk. The weak detectors - Z22, for example - show "
         "the same bimodal picture with the noise peak dominating, and there "
         "this same cut is what digs the real pulses out of the mixture. "
         "(If someone asks what a weak detector looks like, show the backup "
         "slide.)")

# ------------------------------ 6b · fan before vs after the NRMSE cut
s = slide()
title(s, "Aligned fitted curves: before and after the NRMSE cut")
pic(s, "fan_zip7_PBS1_before_all.png", IN(0.3), IN(1.5), IN(6.35), IN(4.6),
    "All fit_ok fitted curves (Z7 PBS1, n = 2008, every curve drawn)")
pic(s, "fan_zip7_PBS1_after_all.png", IN(6.65), IN(1.5), IN(6.35), IN(4.6),
    "After NRMSE ≤ 0.4 (n = 1931): one tight shape family")
notes(s, "The same aligned fitted curves before and after the cut. On the "
         "left, all physical fits - a 200-curve sample is drawn, so the "
         "slow stragglers appear in proportion to their share of the "
         "population. On the right, after the NRMSE cut, a single tight, "
         "consistent shape family remains. That tight family is what makes "
         "a well-defined template possible.")

# ------------------------------------------ 7 · what the cut removes
# ------------------------------------------- 8c · 100-event zoom
s = slide()
title(s, "After NRMSE ≤ 0.4: a closer look")
pic(s, "fan_zip7_PBS1_100zoom.png", IN(1.2), IN(1.5), IN(10.9), IN(5.6),
    "Z7 PBS1: 100 randomly drawn curves of the 1931 kept, zoomed to 25.6\u201326.5 ms")
notes(s, "A closer look at the kept population: one hundred curves randomly "
         "drawn from the 1931 that pass the cut, zoomed to 25 to 27 "
         "milliseconds, so the individual shapes are visible.")

s = slide()
title(s, "The rejected population is noise — verified in the raw traces")
bullets(s, [
    "Example of NRMSE-rejected events (median NRMSE > 0.4 across channels): the raw traces show no pulse",
    "the cut removes noise triggers, not physics",
], IN(0.55), IN(1.2), IN(12.3), IN(1.0), size=14)
def _midlabel(x, w, txt):
    _tb = textbox(s, x, IN(2.12), w, IN(0.4))
    _pp = _tb.text_frame.paragraphs[0]
    _rr = _pp.add_run(); _rr.text = txt
    _rr.font.size, _rr.font.bold, _rr.font.color.rgb, _rr.font.name = Pt(17), True, NAVY, "Arial"
    _pp.alignment = PP_ALIGN.CENTER
_midlabel(IN(0.4), IN(5.9), "Rejected events: raw vs fit")
_midlabel(IN(6.5), IN(6.85), "Aligned curves: kept vs rejected")
pic(s, "slow_rise_zip7_crop.png", IN(0.4), IN(2.75), IN(5.9), IN(4.3),
    "Z7: NRMSE-rejected events — raw traces are noise")
# fan-cut figure enlarged to fill the right column; its legend is pulled out
# into the readable colour-coded text below
# in-figure legend/median text is tiny; the same information is restated at
# normal size right above and below the image
_GREEN, _RED, _BLUE = RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xC0, 0x39, 0x2B), RGBColor(0x4A, 0x70, 0xB0)
_md = textbox(s, IN(6.5), IN(2.48), IN(6.85), IN(0.35))
_pmd = _md.text_frame.paragraphs[0]
for _txt, _col in [("median NRMSE:  ", DARK), ("kept 0.045 (n = 1931)", _GREEN),
                   ("   ·   ", DARK), ("rejected 1.873 (n = 77)", _RED)]:
    _r = _pmd.add_run(); _r.text = _txt
    _r.font.size, _r.font.bold, _r.font.color.rgb, _r.font.name = Pt(14), True, _col, "Arial"
_pmd.alignment = PP_ALIGN.CENTER
pic(s, "fan_cut_zip7_PBS1.png", IN(6.5), IN(2.9), IN(6.85), IN(3.45), stretch=True)
_lg = textbox(s, IN(6.6), IN(6.42), IN(6.7), IN(1.0))
_lg.text_frame.word_wrap = True
_p = _lg.text_frame.paragraphs[0]
for _txt, _col, _bold in [("Z7 PBS1   ", DARK, True),
                          ("── passes the cut (kept)", _GREEN, True),
                          ("      ── cut away = rejected", _RED, True)]:
    _r = _p.add_run(); _r.text = _txt
    _r.font.size, _r.font.bold, _r.font.color.rgb, _r.font.name = Pt(14), _bold, _col, "Arial"
_p2 = _lg.text_frame.add_paragraph()
for _txt, _col, _bold in [("faint blue = measured traces", _BLUE, False),
                          ("      1931 kept / 77 cut (~3.8%)", DARK, False)]:
    _r = _p2.add_run(); _r.text = _txt
    _r.font.size, _r.font.bold, _r.font.color.rgb, _r.font.name = Pt(14), _bold, _col, "Arial"
notes(s, "Before trusting the cut we looked at what it throws away. These are "
         "event grids of the rejected population on Z7: the raw traces show "
         "no pulse at all - they are noise triggers whose slow two-exp fit "
         "happened to converge. The fan-cut view on the right overlays the "
         "passing fitted curves in green and the rejected ones in red on the "
         "aligned data: the green population is the fast physical pulse "
         "shape, the red one is spread out with a median NRMSE forty times "
         "higher - 1.87 against 0.045 - and on this channel only three "
         "point eight percent of the events get cut. So the cut removes "
         "noise triggers, and no real pulses are lost.")

# --------------------------------------- 8 · follow-up 1: slow-fall tail
s = slide()
title(s, "Slow-fall tail in one particular channel")
bullets(s, [
    "PDS2 alone: can not make useful templates; hard to fix in analysis/software, the noise itself needs to be fixed",
    "Temporary solution: use the PDS1 template in place of PDS2 (applied in the ROOT files)",
], IN(0.55), IN(1.25), IN(12.3), IN(1.0), size=14)
pic(s, "slow_fall_zip7_crop.png", IN(0.3), IN(3.5), IN(6.7), IN(3.45),
    "Z7, 3 sampled slow-fall events: normal in PBS2/PCS2/PES2, swings only in PDS2")
# the two t_fall histograms side by side, cropped to the useful 0-7 ms
pic(s, "time_constants_zip7_PAS1_tfall.png", IN(7.05), IN(3.9), IN(3.05), IN(2.2),
    "Z7 PAS1 — normal channel: narrow (median 0.25 ms)")
pic(s, "time_constants_zip7_PDS2_tfall.png", IN(10.2), IN(3.9), IN(3.05), IN(2.2),
    "Z7 PDS2 — bad channel: broad tail (median 0.51 ms)")
_ft = textbox(s, IN(7.05), IN(3.42), IN(6.2), IN(0.5))
_pft = _ft.text_frame.paragraphs[0]
_rft = _pft.add_run(); _rft.text = "Fitted τ_fall distribution, same 0–7 ms axis:"
_rft.font.size, _rft.font.bold, _rft.font.color.rgb, _rft.font.name = Pt(15), True, NAVY, "Arial"
notes(s, "The first lead comes from the fan plot after the cut: some curves "
         "still fall very slowly. We chased it by sampling: among events "
         "passing the cut, take fall times above one point five "
         "milliseconds, randomly sample ten, and draw raw versus fit in "
         "every channel. Two findings. First, the sampled events are real "
         "pulses - so they stay in, and we apply no fall cut. Second, their "
         "extreme fall times all trace back to one channel: only PDS2 is "
         "swinging wildly, while the same events are normal fast pulses in "
         "the other eleven channels - so the extreme tail is a one-channel "
         "low-frequency artifact, not a genuinely slow pulse. PDS2 alone "
         "cannot make a useful template, and this is hard to fix on the "
         "analysis or software side - the noise itself needs to be fixed in "
         "hardware. As a temporary solution we use the PDS1 template in "
         "place of PDS2, and that substitution is applied in the delivered "
         "ROOT files.")

# ------------------------------------------------- 11 · template family 1
s = slide()
title(s, "Results: 1x1 Templates")
bullets(s, [
    "NRMSE-weighted mean of the fit_ok 2-exp curves",
    "Smooth & noise-free by construction",
], IN(0.55), IN(1.25), IN(12.3), IN(1.1), size=18)
pic(s, "template_overlay_zip7_PBS1.png", IN(0.5), IN(2.5), IN(12.3), IN(2.15),
    "Z7 PBS1, blue: the fitted curves (200 of 1931 drawn);  red: the 1x1 template")
pic(s, "mean_compare_zip7_PBS1.png", IN(0.7), IN(4.95), IN(11.9), IN(2.25),
    "Solid red: NRMSE-weighted mean (= the 1x1 template);  dashed navy: plain mean of the PCA input curves (= nxm0)")
notes(s, "First template family: the analytic one. For each channel we take "
         "every physical fitted curve at the common pretrigger and average them "
         "with a weight of one over NRMSE squared - badly fit events count "
         "less but nothing is excluded by hand. After the cuts, shown here, "
         "a single tight shape family remains, and its weighted mean is the "
         "1x1 template - smooth and noise-free by construction because it is "
         "built from analytic curves, delivered as peak-normalized ROOT "
         "histograms.")

# ------------------------------------------------- 12 · template family 2
s = slide()
title(s, "Results: NxM Templates")
bullets(s, [
    "channel-specific PCA over the fitted curves (fit_ok + NRMSE ≤ 0.4)",
    "nxm0 = mean shape;  nxm1–4 = principal components;  real pulse = Σᵢ ampᵢ · nxmᵢ",
    "PC1 + PC2 ≈ 96–98% of the shape variance",
], IN(0.55), IN(1.22), IN(12.3), IN(1.5), size=16)
pic(s, "nxm_plain_zip7_PBS1.png", IN(2.4), IN(2.85), IN(8.5), IN(4.3),
    "Z7 PBS1: nxm0 (plain mean over the fit_ok + NRMSE-cut curves) + nxm1–4, zoomed to the pulse")
notes(s, "Second family: the NxM PCA templates, built to capture the "
         "pulse-shape variation across events. We run a PCA over the fitted curves that "
         "pass all cuts, in a window around the pulse. The mean shape "
         "becomes template zero, and the first four principal components "
         "become templates one to four - they oscillate and can be "
         "negative, and the optimal filter fits a real pulse as a linear "
         "combination of them. The first two components already capture 96 "
         "to 98 percent of the shape variance on every detector.")

# ------------------------------------------------- 13 · summary
# ------------------------------- 13b · NxM processing example (UMN run)
s = slide()
title(s, "NxM processing example: raw vs reconstruction")
pic(s, "nxm_reco_overlay_s1.png", IN(0.4), IN(1.3), IN(12.5), IN(2.8))
pic(s, "nxm_reco_overlay_s2.png", IN(0.4), IN(4.15), IN(12.5), IN(2.8),
    "Z7, 3 example K-line events (top: S1 channels, bottom: S2): raw (gray) / LP (blue) vs Σ ampₖ·nxmₖ (red), amplitudes from the UMN (Addison) processing of our templates")
notes(s, "A first look at the NxM processing exercised on our delivered "
         "templates: UMN's run by Addison produced the five template "
         "amplitudes per channel, and the red curve is the reconstruction, "
         "the sum of amplitude k times template k, overlaid on the raw "
         "trace. On the S1 channels the reconstruction follows the pulse "
         "shape closely; some S2 channels show large oscillating "
         "components, including PDS2 where the substituted PDS1 template "
         "is in use - exactly the kind of feedback this exercise is for.")

s = slide()
title(s, "Template file")
bullets(s, [
    "Two template sets for all 13 detectors: 1x1 (2-exp weighted) + NxM PCA",
    "In cdmsbats_config feature branch, ready for merge",
], IN(0.55), IN(1.75), IN(12.3), IN(1.4), size=22)
_fh = textbox(s, IN(0.55), IN(3.55), IN(12.3), IN(0.6))
_pf = _fh.text_frame.paragraphs[0]
_rf = _pf.add_run(); _rf.text = "Future Steps"
_rf.font.size, _rf.font.bold, _rf.font.color.rgb, _rf.font.name = Pt(26), True, NAVY, "Arial"
bullets(s, [
    "Three-, four-exponential fits",
    "Exercise NxM processing",
    "Other detectors: worse noise, a few more cuts added, in the backup slides",
], IN(0.55), IN(4.25), IN(12.3), IN(2.2), size=22)
notes(s, "Both template sets are ready for all 13 detectors: the analytic "
         "1x1 templates and the PCA NxM set, in the official PulseTemplates "
         "format, sitting in the cdmsbats underscore config feature branch "
         "and ready for merge. Future steps: extend the fit to three and "
         "four exponentials, exercise the NxM processing chain on the new "
         "templates, and the other detectors - they have worse noise and "
         "needed a few more cuts, all documented in the backup slides. If "
         "anyone is interested, I am happy to walk through the backup "
         "slides. Thank you - happy to take questions.")

# ------------------------------------------- backup · weak detectors (Z22)
s = slide()
title(s, "Backup · weak detectors (example: Z22)")
bullets(s, [
    "Z1 / Z4 / Z6 / Z18 / Z19 / Z22 / Z24: K-line inside the noise population, the window admits a noise-dominated mixture",
    "Same bimodal NRMSE picture, noise peak dominant; the 0.4 cut digs the real pulses out (Z22 PCS1: 7198 kept / 4788 cut, ~40%)",
    "Kept bundle broader; steep red curves = fits latching onto noise spikes, no real pulses lost",
    "τ_rise ≤ 0.3 ms (after NRMSE): removes 55–74% on weak zips (Z22: 71%) vs 1.6% on Z7; removes residual slow drift",
], IN(0.55), IN(1.1), IN(12.3), IN(2.15), size=13)
pic(s, "nrmse_zip22_PAS1.png", IN(0.4), IN(3.35), IN(6.3), IN(1.7),
    "Z22 PAS1 NRMSE histogram: noise dominates")
pic(s, "fan_cut_zip22_PCS1.png", IN(6.9), IN(3.35), IN(6.3), IN(1.7),
    "Z22 PCS1: kept (green) vs cut (red)")
pic(s, "fan_final_zip22_PCS1.png", IN(2.2), IN(5.55), IN(8.9), IN(1.55),
    "Z22 PCS1 after NRMSE ≤ 0.4 + τ_rise ≤ 0.3 ms: final shape family (760 fits)")
notes(s, "Backup. On the weak detectors the K-line overlaps the noise "
         "population, so the PTOF window admits a noise-dominated mixture. "
         "The NRMSE distribution is still bimodal but the noise peak "
         "dominates, and the same 0.4 cut is what extracts the real pulses "
         "- on Z22 PCS1 about forty percent of the physical fits are "
         "removed. The kept bundle is broader than on a quiet detector, and "
         "the steep red curves are fits latching onto sharp noise spikes, "
         "not fast pulses being thrown away. The rise-time ceiling then "
         "removes fifty-five to seventy-four percent of what survived the "
         "NRMSE cut on the weak zips - residual slow drift - versus one "
         "point six percent on Z7. The surviving shape family is at the "
         "bottom.")

# ------------------------------------------- backup · τ_rise ≤ 0.3 ms ceiling
s = slide()
title(s, "Backup · τ_rise ≤ 0.3 ms ceiling (PCA input)")
bullets(s, [
    "Slow baseline drift survives NRMSE (a slow 2-exp hugs it) → mimics a slow rise",
    "Fast pulses at τ_rise ≈ 0.1 ms (p90 ≈ 0.15 ms) → ceiling at 0.3 ms blocks the drift tail",
    "Removed after the NRMSE cut: Z7 1.6% (quiet 2–6%), weak 55–74% (Z22 71%), all zips 54%",
    "Trade-off: trims the very slowest genuine pulses → decision pending",
], IN(0.55), IN(1.2), IN(12.3), IN(2.15), size=14)
pic(s, "time_constants_zip7_PAS1.png", IN(1.4), IN(3.6), IN(10.5), IN(2.7),
    "Z7 PAS1: fitted τ_rise (median 0.105 ms) and τ_fall distributions")
notes(s, "Backup for the rise-time cut. On weak-window detectors a slow "
         "baseline drift also fits as a slow rise, with a tiny residual, so "
         "NRMSE cannot catch it. Real fast pulses cluster near 0.1 ms while "
         "the drift tail stretches much further, so the PCA input gets a "
         "ceiling at 0.3 ms. Relative to the NRMSE step it removes almost "
         "nothing on quiet detectors - 1.6% on Z7, mostly the bad channel "
         "PDS2 - but 55 to 74% on the weak ones, 71% on Z22, pooled 54% over "
         "all zips. The price is that the very slowest genuine pulses get "
         "trimmed too; that trade-off is documented and not final.")

# ---------------------------- backup · τ_rise cut, good vs bad channel (Z7)
s = slide()
title(s, "Backup · τ_rise cut: good vs bad channel (Z7)")
bullets(s, [
    "Good channel PBS1: tight clean bundle; the cut removes ~0% (nearly free on quiet channels)",
    "Bad channel PDS2: broad messy fan, the channel itself is bad; the cut trims slow risers, 1462 → 1154 (21%)",
], IN(0.55), IN(1.15), IN(12.3), IN(1.35), size=14)
pic(s, "fan_zip7_PBS1_nrmse.png", IN(0.35), IN(2.7), IN(6.25), IN(1.75),
    "PBS1 (good): NRMSE ≤ 0.4")
pic(s, "fan_zip7_PBS1_after.png", IN(6.85), IN(2.7), IN(6.25), IN(1.75),
    "PBS1 (good): + τ_rise ≤ 0.3 ms (unchanged)")
pic(s, "fan_zip7_PDS2_nrmse.png", IN(0.35), IN(5.05), IN(6.25), IN(1.75),
    "PDS2 (bad): NRMSE ≤ 0.4")
pic(s, "fan_zip7_PDS2_after.png", IN(6.85), IN(5.05), IN(6.25), IN(1.75),
    "PDS2 (bad): + τ_rise ≤ 0.3 ms (slow risers trimmed)")
notes(s, "Backup showing the cut channel-by-channel on Z7. On the good "
         "channel PBS1 the fitted curves are already a tight, clean bundle, "
         "and adding the rise-time ceiling changes essentially nothing - it "
         "removes about zero percent. On the bad channel PDS2 the fan is "
         "broad and messy, which by itself shows the channel is misbehaving; "
         "there the ceiling trims the slow-rising curves, 1462 down to 1154, "
         "about twenty-one percent. So the cut is nearly free where the data "
         "is clean and does real work only where a channel is bad.")

# ---------------------------- backup · what the τ_rise cut removes (Z22)
s = slide()
title(s, "Backup · what the τ_rise cut removes (Z22)")
bullets(s, [
    "Events that PASS NRMSE ≤ 0.4 but are REMOVED by τ_rise ≤ 0.3 ms: the noise 0.4 let through, 0.3 catches",
    "PCS1 / PDS1 pass at NRMSE ≈ 0.15 (a slow 2-exp hugs the drift); the raw trace has no clean fast pulse",
    "Trade-off: also trims genuine slow-rise pulses; decision pending",
], IN(0.55), IN(1.15), IN(12.3), IN(1.7), size=14)
pic(s, "drift_zip22_crop.png", IN(2.1), IN(3.25), IN(9.1), IN(3.8),
    "Z22, 3 events × PAS1/PBS1/PCS1/PDS1: pass NRMSE ≤ 0.4, removed by τ_rise ≤ 0.3 ms; PCS1/PDS1 = pure slow drift")
notes(s, "Backup: what the rise-time cut actually removes. These are exactly "
         "the Z22 events that pass the NRMSE cut - median NRMSE below 0.4 - "
         "but are removed by the rise-time ceiling, median t_rise above 0.3 "
         "ms. In other words, the noise that 0.4 let through and 0.3 catches. "
         "Look at PCS1 and PDS1: their NRMSE is low, around 0.15, because a "
         "slow two-exp hugs the trace with a tiny residual - but the raw "
         "trace is just a slow baseline drift, there is no clean fast pulse. "
         "The honest caveat is that the same cut also trims a small number of "
         "genuine slow-rise pulses, which is the documented, still-open "
         "trade-off.")



_SPLIT_DIR = os.path.join(FIG, "backup_zip7", "_split")
os.makedirs(_SPLIT_DIR, exist_ok=True)


def _white_rows(img):
    """Boolean list: row is (almost) pure background white."""
    import numpy as np
    g = np.asarray(img.convert("L"))
    return (g.min(axis=1) > 247)


def snap_cut(white, target, window):
    """Nearest all-white row to target within +-window (else target)."""
    lo = max(1, target - window)
    hi = min(len(white) - 1, target + window)
    best, bestd = None, None
    for y in range(lo, hi):
        if white[y]:
            d = abs(y - target)
            if bestd is None or d < bestd:
                best, bestd = y, d
    return best if best is not None else target


def _white_cols(img):
    """Boolean list: column is (almost) pure background white."""
    import numpy as np
    g = np.asarray(img.convert("L"))
    return (g.min(axis=0) > 247)


def _content_blocks(mask):
    """Contiguous non-white index ranges."""
    blocks, start = [], None
    for i, w in enumerate(mask):
        if not w and start is None:
            start = i
        elif w and start is not None:
            blocks.append((start, i)); start = None
    if start is not None:
        blocks.append((start, len(mask)))
    return blocks


def _panel_units(img):
    """Vertical (start, end) unit per panel: the panel box plus its title
    above and tick/label text below. Panels are the tall content blocks
    (inside a plot box no row is all-white thanks to the frame spines);
    text lines are the short ones. Header text above the first panel is
    dropped — the slide carries the description."""
    white = _white_rows(img)
    blocks = _content_blocks(white)
    hmax = max(b[1] - b[0] for b in blocks)
    thr = max(60, int(hmax * 0.35))
    panels = [b for b in blocks if b[1] - b[0] >= thr]
    small = [b for b in blocks if b[1] - b[0] < thr]
    starts = []
    for ps, pe in panels:
        above = sorted(b for b in small if b[1] <= ps + 2)
        top = ps
        if above and ps - above[-1][1] <= 12:      # the title hugs its panel
            top = above[-1][0]
            for b in reversed(above[:-1]):         # merge artifact slivers
                if top - b[1] <= 6:
                    top = b[0]
                else:
                    break
        starts.append(max(0, top - 4))
    units, tails = [], []
    for i, (ps, pe) in enumerate(panels[:-1]):
        end = starts[i + 1] - 4
        units.append((starts[i], end))
        tails.append(end - pe)
    ps, pe = panels[-1]
    tail = sorted(tails)[len(tails) // 2] if tails else 60
    units.append((starts[-1], min(img.size[1], pe + tail)))
    return units


GAL = "backup_zip7"


def _gal_slide(head, sub_txt, info):
    s = slide()
    title(s, head, sub=sub_txt)
    tb = textbox(s, IN(0.55), IN(1.18), IN(12.3), IN(0.62))
    _p = tb.text_frame.paragraphs[0]
    _r = _p.add_run(); _r.text = info
    _r.font.size, _r.font.color.rgb, _r.font.name = Pt(11.5), DARK, "Arial"
    return s


def _place_stack(s, crops, l, t, box_w, box_h, gap=Inches(0.22)):
    """Stack crops vertically, each aspect-fit in an equal cell, centered."""
    n = len(crops)
    cell_h = (box_h - gap * (n - 1)) / n
    for i, (cp, (cw, ch)) in enumerate(crops):
        sc = min(box_w / cw, cell_h / ch)
        dw, dh = int(cw * sc), int(ch * sc)
        left = l + int((box_w - dw) / 2)
        top = t + i * (cell_h + gap) + int((cell_h - dh) / 2)
        s.shapes.add_picture(cp, left, top, width=dw, height=dh)


def channel_fig_pages(fname, short, info, per_page=2):
    """One channel-stacked figure -> pages of per_page full-width strips."""
    img = Image.open(os.path.join(FIG, GAL, fname))
    w0 = img.size[0]
    units = _panel_units(img)
    n = len(units)
    base = os.path.splitext(fname)[0]
    npages = (n + per_page - 1) // per_page
    for p in range(npages):
        i0, i1 = p * per_page, min((p + 1) * per_page, n)
        crops = []
        for i in range(i0, i1):
            crop = img.crop((0, units[i][0], w0, units[i][1]))
            cp = os.path.join(_SPLIT_DIR, f"{base}_r{i+1}.png")
            crop.save(cp)
            crops.append((cp, crop.size))
        s = _gal_slide(f"Backup · Z7 · {short}  ({p+1}/{npages})",
                       f"figure: {fname}, channel panels {i0+1}–{i1} of {n}, "
                       f"top to bottom",
                       info)
        _place_stack(s, crops, IN(0.55), IN(1.95), IN(12.3), IN(5.3))
        notes(s, f"Backup gallery, Z7, {short}, page {p+1} of {npages}. {info}")


def grid_fig_pages(fname, short, info, rows_per_page=3):
    """Event x channel grid -> pages of rows_per_page event rows, channels
    split into a left and a right half so every panel stays readable."""
    img = Image.open(os.path.join(FIG, GAL, fname))
    w0 = img.size[0]
    units = _panel_units(img)
    n = len(units)
    body = img.crop((0, units[0][0], w0, units[-1][1]))
    xh = snap_cut(_white_cols(body), w0 // 2, w0 // 10)
    base = os.path.splitext(fname)[0]
    groups = [(a, min(a + rows_per_page, n))
              for a in range(0, n, rows_per_page)]
    pages = [(a, b, side) for a, b in groups for side in (0, 1)]
    for p, (a, b, side) in enumerate(pages):
        x0, x1 = (0, xh) if side == 0 else (xh, w0)
        half = "left" if side == 0 else "right"
        crop = img.crop((x0, units[a][0], x1, units[b - 1][1]))
        cp = os.path.join(_SPLIT_DIR, f"{base}_e{a+1}-{b}_{half}.png")
        crop.save(cp)
        s = _gal_slide(f"Backup · Z7 · {short}  ({p+1}/{len(pages)})",
                       f"figure: {fname}, event rows {a+1}–{b} of {n}, "
                       f"{half} half of the channels",
                       info)
        _place_stack(s, [(cp, crop.size)], IN(0.55), IN(1.95), IN(12.3), IN(5.3))
        notes(s, f"Backup gallery, Z7, {short}, page {p+1} of {len(pages)}. {info}")


# ---------------------------- backup · fit_ok-rejected events (late pulses)
s = slide()
title(s, "Backup · fit_ok-rejected events: normal amplitude, late pulse")
pic(s, "zip7_fitok_rejected_events_3x4.png", IN(0.8), IN(1.35), IN(11.7), IN(5.75),
    "Z7, 3 fit_ok-rejected events × PAS1/PBS1/PCS1/PFS1, full 0–52 ms trace: the real pulse sits far right of the fit window; red (!ok) fits are the negative solutions")
notes(s, "Backup: what fit_ok actually rejects. Sampled events with the full "
         "52 millisecond trace drawn: the amplitude is perfectly normal "
         "K-line, but the pulse arrives tens of milliseconds late, outside "
         "the pretrigger search range of the fit, so the fit only sees noise "
         "and lands on the swapped-tau negative solution. Removing these is "
         "correct because such pulses cannot enter a template anyway.")

# ------------------------------------------------ gallery intro / contents
s = slide()
title(s, "Backup · Z7 full results gallery",
      sub="every diagnostic and result figure of the Z7 pipeline, in processing order")
bullets(s, [
    "Dataset: the K-line selection window",
    "Example grids: LP + fit and raw vs fit, event × channel",
    "Fit diagnostics: NRMSE, fitted pretrigger, τ_rise / τ_fall distributions",
    "Aligned fans: no cut → NRMSE ≤ 0.4 → + τ_rise ≤ 0.3 ms; kept vs rejected; measured traces",
    "Special populations: rejected events, well-fit slow-rise events, PDS2 slow-fall study",
    "Templates: PCA nxm0–4 per channel, summed PT / PS1 / PS2",
    "Tall figures are cut at white space into full-width panels, two per page; grids into 3-event × half-channel blocks",
], IN(0.55), IN(1.6), IN(12.3), IN(4.6), size=16)
notes(s, "The full Z7 gallery: every figure the pipeline produces, in "
         "processing order, split into readable pieces.")

# ------------------------------------------------ dataset window
s = _gal_slide("Backup · Z7 · K-line selection window",
               "figure: kline_zip7.png",
               "Summed PTOFamps spectrum of Z7; red line = Prof. Saab's 10.37 keV K-line position, "
               "shaded band = the ×/÷ 1.35 selection window; everything inside is cached raw.")
pic(s, "kline_zip7.png", IN(2.65), IN(2.0), IN(8.0), IN(5.2),
    caption="Z7 PTOFamps spectrum with the K-line selection window")
notes(s, "Backup gallery: the Z7 K-line window on the PTOFamps spectrum.")

# ------------------------------------------------ example grids
grid_fig_pages("zip7_fit_examples.png",
    "fit_examples: LP trace vs 2-exp fit",
    "First 15 cached events × all channels; blue = 100 kHz LP trace, red = 2-exp fit. "
    "A real event fits well in every channel at once; a noise trigger fails across the whole row.")
grid_fig_pages("zip7_raw_vs_fit_examples.png",
    "raw_vs_fit_examples: raw trace vs fit",
    "Same 15 events × all channels; gray = raw unfiltered trace, blue = LP, red = fit, "
    "per-panel NRMSE top-right. Noise vs pulse size is directly visible in the raw data.")

# ------------------------------------------------ fit diagnostics
channel_fig_pages("zip7_nrmse.png",
    "nrmse: NRMSE distribution per channel",
    "NRMSE = RMS(fit residual)/fitted peak, log-log. Bimodal: good fits (~0.05–0.1) vs "
    "noise triggers (~1–2); the valley at ≈0.4 sets quality cut 2.")
channel_fig_pages("zip7_pretrigger.png",
    "pretrigger: fitted onset per channel",
    "The onset is a FREE fit parameter; fitted values cluster ≈230 samples after the nominal "
    "16050; pinning it inside the fit would bias every other parameter.")
channel_fig_pages("zip7_time_constants.png",
    "time_constants: τ_rise / τ_fall per channel",
    "Fast-pulse population at τ_rise ≈ 0.1 ms; isolated spikes at the parameter bounds are "
    "noise fits pinned at the fit limits, removed by the NRMSE cut.")

# ------------------------------------------------ aligned fans
channel_fig_pages("zip7_fitted_curves_overlay.png",
    "fitted_curves_overlay: fan, no cut",
    "Every fit_ok fitted curve at the common pretrigger 16050, peak-normalized. Pure shape "
    "distribution; the broad slow curves are noise fits that the cuts remove.")
channel_fig_pages("zip7_fitted_curves_overlay_nrmse0.4.png",
    "fitted_curves_overlay: after NRMSE ≤ 0.4",
    "Same fan after the NRMSE cut: the noise fan collapses, the fast-pulse bundle remains.")
channel_fig_pages("zip7_fitted_curves_overlay_nrmse0.4_trise0.30ms.png",
    "fitted_curves_overlay: after both cuts",
    "After NRMSE ≤ 0.4 AND τ_rise ≤ 0.3 ms, the exact PCA input population: "
    "the clean fast-pulse bundle only.")
channel_fig_pages("zip7_overlay_fan_cut_nrmse0.4.png",
    "overlay_fan_cut: kept vs rejected fits",
    "Gray-blue = aligned measured traces; green = fits with NRMSE ≤ 0.4 (kept); red = "
    "NRMSE > 0.4 (cut away). Median NRMSE and counts of both populations stamped per channel.")
channel_fig_pages("zip7_lp_aligned_overlay.png",
    "aligned_overlay: aligned measured traces",
    "Blue = measured LP traces shifted by (fitted pretrigger − 16050); red = mean of all "
    "fit_ok events; orange = NRMSE-weighted mean of the fitted curves. Rise edges line up at 16050.")

# ------------------------------------------------ special populations
grid_fig_pages("zip7_slow_rise_events.png",
    "slow_rise_events: NRMSE-rejected events",
    "Events with median NRMSE > 0.4 across channels: the raw traces show no pulse in any "
    "channel: noise triggers, not physics.")
grid_fig_pages("zip7_shadow_events.png",
    "shadow_events: well-fit slow-rise events",
    "Median NRMSE ≤ 0.4 AND median τ_rise > 0.2 ms: genuinely slow, well-fit pulses; kept, "
    "and exactly what the NxM multi-templates are for.")
grid_fig_pages("zip7_PDS2_slow_fall.png",
    "slow_fall_events: long-τ_fall study (PDS2)",
    "10 random events with τ_fall > 1.5 ms drawn in every channel: 11 channels show normal "
    "fast pulses; the long fall is the PDS2 low-frequency hardware noise, not slow physics.")

# ------------------------------------------------ templates
channel_fig_pages("zip7_pca_templates.png",
    "deliverables/nxm: PCA templates nxm0–4",
    "nxm0 = mean curve; nxm1–4 = PCA components (oscillating basis vectors, may be negative), "
    "all peak-normalized; PC1+PC2 capture 96–98% of the shape variance. PDS2 replaced by the "
    "PDS1 template in the delivered files.")

# summed 1x1 templates: three figures on one slide
s = slide()
title(s, "Backup · Z7 · 1x1 summed templates PT / PS1 / PS2",
      sub="figures: deliverables/1x1/plots/{PT,PS1,PS2}/zip7_*.png")
tb = textbox(s, IN(0.55), IN(1.18), IN(12.3), IN(0.6))
_p = tb.text_frame.paragraphs[0]
_r = _p.add_run()
_r.text = ("Single-template (1x1) summed curves: peak-normalized average of the per-channel nxm0 templates: "
           "PT = all channels, PS1 / PS2 = side-1 / side-2 only.")
_r.font.size, _r.font.color.rgb, _r.font.name = Pt(12.5), DARK, "Arial"
for _k, _f in enumerate(["zip7_PT.png", "zip7_PS1.png", "zip7_PS2.png"]):
    _img = Image.open(os.path.join(FIG, GAL, _f))
    _u = _panel_units(_img)
    _crop = _img.crop((0, _u[0][0], _img.size[0], _u[-1][1]))
    _cp = os.path.join(_SPLIT_DIR, _f.replace(".png", "_clean.png"))
    _crop.save(_cp)
    pic(s, os.path.join(GAL, "_split", _f.replace(".png", "_clean.png")),
        IN(0.4 + 4.3 * _k), IN(2.3), IN(4.1), IN(4.2),
        caption=_f.replace("zip7_", "Z7 ").replace(".png", ""))
notes(s, "Backup gallery: the three summed 1x1 templates of Z7 — PT over all "
         "channels, PS1 and PS2 over each side, all peak-normalized averages "
         "of the channel nxm0 templates.")

# page numbers: "N / total" bottom-right, skip the title slide
_slides = list(prs.slides)
_total = len(_slides)
for _i, _s in enumerate(_slides, start=1):
    if _i == 1:
        continue
    _pn = _s.shapes.add_textbox(SW - Inches(1.35), SH - Inches(0.45),
                                Inches(1.1), Inches(0.3))
    _p = _pn.text_frame.paragraphs[0]
    _r = _p.add_run(); _r.text = f"{_i} / {_total}"
    _r.font.size, _r.font.color.rgb, _r.font.name = Pt(11), GRAY, "Arial"
    _p.alignment = PP_ALIGN.RIGHT

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
